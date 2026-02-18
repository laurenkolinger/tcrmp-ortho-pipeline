#!/usr/bin/env python3
"""
Generate a PowerPoint from the current project's edited TIFs.

Reads .current_project to find data/{PROJECT_DIR}/edited/, then builds
a PPTX with 3 images per slide, grouped by site.

Usage:
    python3 src/create_ortho_pptx.py
    python3 src/create_ortho_pptx.py -o output/custom_name.pptx
"""

import argparse
import os
import re
import subprocess
import tempfile
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

PROJECT_FILE = ".current_project"


def natural_sort_key(t):
    return [int(p) for p in re.findall(r'\d+', t)]


def discover_files(src_dir):
    """Read {SITE}_{TRANSECT}_full.tif files from a directory."""
    site_files = defaultdict(list)
    for fname in sorted(os.listdir(src_dir)):
        if not fname.endswith("_full.tif"):
            continue
        match = re.match(r'^([A-Z]{3})_(T\d+(?:_\d+)?)_full\.tif$', fname)
        if not match:
            continue
        site_code = match.group(1)
        transect = match.group(2)
        tif_path = os.path.join(src_dir, fname)
        site_files[site_code].append((transect, tif_path))

    for site in site_files:
        site_files[site].sort(key=lambda x: natural_sort_key(x[0]))
    return site_files


def tif_to_jpeg(tif_path, tmp_dir):
    basename = os.path.splitext(os.path.basename(tif_path))[0]
    jpg_path = os.path.join(tmp_dir, basename + ".jpg")
    subprocess.run(
        ["sips", "-s", "format", "jpeg", "-s", "formatOptions", "80",
         tif_path, "--out", jpg_path],
        capture_output=True, check=True
    )
    return jpg_path


def main():
    parser = argparse.ArgumentParser(
        description="Generate a PPTX from the current project's edited TIFs."
    )
    parser.add_argument(
        "--output", "-o",
        default=None,
        help="Output PPTX path (default: output/TCRMP_{PROJECT_DIR}.pptx)"
    )
    args = parser.parse_args()

    if not os.path.isfile(PROJECT_FILE):
        print("Error: .current_project not found.")
        print("Run copy_orthomosaics.py first to set the project.")
        return 1

    with open(PROJECT_FILE) as f:
        project_dir = f.read().strip()

    src_dir = os.path.join("data", project_dir, "edited")
    if not os.path.isdir(src_dir):
        print(f"Error: {src_dir} not found.")
        return 1

    output_pptx = args.output or os.path.join("output", f"TCRMP_{project_dir}.pptx")
    os.makedirs(os.path.dirname(output_pptx) or ".", exist_ok=True)

    site_files = discover_files(src_dir)
    sorted_sites = sorted(site_files.keys())

    total = sum(len(v) for v in site_files.values())
    print(f"Found {len(sorted_sites)} sites, {total} images in {src_dir}\n")
    for site in sorted_sites:
        files = site_files[site]
        print(f"  {site}: {', '.join(t for t, _ in files)}")
    print()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    with tempfile.TemporaryDirectory() as tmp_dir:
        for site in sorted_sites:
            files = site_files[site]

            for page_idx in range(0, len(files), 3):
                chunk = files[page_idx:page_idx + 3]
                n = len(chunk)

                slide = prs.slides.add_slide(prs.slide_layouts[6])

                txBox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.1), Inches(12), Inches(0.5)
                )
                p = txBox.text_frame.paragraphs[0]
                suffix = " (continued)" if page_idx > 0 else ""
                p.text = f"{site}{suffix}"
                p.font.size = Pt(28)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

                top_margin = Inches(0.7)
                available_height = Inches(6.5)
                img_height = available_height / 3 - Inches(0.3)

                for i, (transect, tif_path) in enumerate(chunk):
                    print(f"  Processing {site} {transect}...")
                    jpg_path = tif_to_jpeg(tif_path, tmp_dir)

                    with Image.open(jpg_path) as img:
                        img_w, img_h = img.size

                    aspect = img_w / img_h
                    slot_top = top_margin + i * (img_height + Inches(0.3))

                    max_width = Inches(12.5)
                    calc_width = img_height * aspect
                    if calc_width > max_width:
                        final_width = max_width
                        final_height = max_width / aspect
                    else:
                        final_width = calc_width
                        final_height = img_height

                    left = (prs.slide_width - final_width) / 2

                    label_box = slide.shapes.add_textbox(
                        Inches(0.5), slot_top, Inches(12), Inches(0.25)
                    )
                    label_p = label_box.text_frame.paragraphs[0]
                    label_p.text = transect
                    label_p.font.size = Pt(14)
                    label_p.font.bold = True
                    label_p.alignment = PP_ALIGN.CENTER

                    img_top = slot_top + Inches(0.25)
                    slide.shapes.add_picture(
                        jpg_path, int(left), int(img_top),
                        int(final_width), int(final_height)
                    )

                print(f"  Slide for {site} (images {page_idx+1}-{page_idx+n})")

    prs.save(output_pptx)
    print(f"\nSaved: {output_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
